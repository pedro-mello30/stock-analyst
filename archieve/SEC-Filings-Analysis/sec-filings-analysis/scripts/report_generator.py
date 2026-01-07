#!/usr/bin/env python3
"""
Report Generator for SEC Filings Analysis

Generates professional PPTX and DOCX reports with financial analysis,
charts, and investment theses templates.
"""

import os
import json
from typing import Dict, List, Any, Optional, Union
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Inches as DocxInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import logging

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class ReportGenerator:
    """Main class for generating professional analysis reports."""

    def __init__(self, template_dir: str = "assets/templates"):
        """
        Initialize report generator.

        Args:
            template_dir: Directory containing report templates
        """
        self.template_dir = template_dir
        self.output_dir = "reports"

        # Ensure output directory exists
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

    def generate_analysis_report(self, company_data: Dict[str, Any],
                               charts: Dict[str, str],
                               risk_assessment: Dict[str, Any],
                               company_name: str = "Company") -> str:
        """
        Generate detailed analysis report (PPTX format).

        Args:
            company_data: Complete financial analysis data
            charts: Dictionary of chart file paths
            risk_assessment: Risk assessment results
            company_name: Company name for report

        Returns:
            Path to generated report
        """
        # Create presentation
        prs = Presentation()

        # Generate report slides
        self._add_title_slide(prs, company_name)
        self._add_executive_summary_slide(prs, company_data, risk_assessment)
        self._add_financial_highlights_slide(prs, company_data)
        self._add_ratio_analysis_slides(prs, company_data, charts)
        self._add_trend_analysis_slides(prs, company_data, charts)
        self._add_risk_assessment_slide(prs, risk_assessment)
        self._add_valuation_metrics_slide(prs, company_data)
        self._add_investment_recommendations_slide(prs, company_data, risk_assessment)

        # Save presentation
        filename = f"{company_name.replace(' ', '_')}_Analysis_Report.pptx"
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)

        logger.info(f"Detailed analysis report saved: {output_path}")
        return output_path

    def generate_investment_thesis(self, company_data: Dict[str, Any],
                                 thesis_framework: Dict[str, Any],
                                 company_name: str = "Company") -> str:
        """
        Generate investment thesis document (DOCX format).

        Args:
            company_data: Complete financial analysis data
            thesis_framework: Investment thesis framework
            company_name: Company name for report

        Returns:
            Path to generated report
        """
        doc = Document()

        # Add title page
        self._add_title_page(doc, company_name, "Investment Thesis")

        # Generate thesis sections
        self._add_company_overview_section(doc, company_data)
        self._add_industry_analysis_section(doc, thesis_framework.get('industry', {}))
        self._add_financial_analysis_section(doc, company_data)
        self._add_competitive_analysis_section(doc, thesis_framework.get('competitive', {}))
        self._add_risk_assessment_section(doc, thesis_framework.get('risks', {}))
        self._add_valuation_section(doc, company_data, thesis_framework.get('valuation', {}))
        self._add_investment_thesis_section(doc, thesis_framework.get('thesis', {}))
        self._add_conclusion_section(doc, thesis_framework.get('conclusion', {}))

        # Save document
        filename = f"{company_name.replace(' ', '_')}_Investment_Thesis.docx"
        output_path = os.path.join(self.output_dir, filename)
        doc.save(output_path)

        logger.info(f"Investment thesis saved: {output_path}")
        return output_path

    def generate_executive_summary(self, key_metrics: Dict[str, Any],
                                 company_name: str = "Company") -> str:
        """
        Generate executive summary report (PPTX format).

        Args:
            key_metrics: Key financial metrics and highlights
            company_name: Company name for report

        Returns:
            Path to generated summary
        """
        prs = Presentation()

        # Create concise executive summary slides
        self._add_title_slide(prs, company_name, "Executive Summary")
        self._add_key_metrics_slide(prs, key_metrics)
        self._add_investment_highlights_slide(prs, key_metrics)
        self._add_risk_summary_slide(prs, key_metrics)

        # Save presentation
        filename = f"{company_name.replace(' ', '_')}_Executive_Summary.pptx"
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)

        logger.info(f"Executive summary saved: {output_path}")
        return output_path

    def populate_template(self, template_path: str, data: Dict[str, Any],
                        charts: Dict[str, str]) -> str:
        """
        Populate existing template with data and charts.

        Args:
            template_path: Path to template file
            data: Data to populate template
            charts: Chart file paths to insert

        Returns:
            Path to populated template
        """
        if template_path.endswith('.pptx'):
            return self._populate_pptx_template(template_path, data, charts)
        elif template_path.endswith('.docx'):
            return self._populate_docx_template(template_path, data, charts)
        else:
            raise ValueError("Unsupported template format. Use .pptx or .docx")

    def export_reports(self, output_dir: str, company_name: str) -> Dict[str, str]:
        """
        Export all generated reports to specified directory.

        Args:
            output_dir: Directory to export reports
            company_name: Company name for file naming

        Returns:
            Dictionary of exported report paths
        """
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        exported_reports = {}

        # Copy all reports to export directory
        for filename in os.listdir(self.output_dir):
            if filename.startswith(company_name.replace(' ', '_')):
                src_path = os.path.join(self.output_dir, filename)
                dst_path = os.path.join(output_dir, filename)
                import shutil
                shutil.copy2(src_path, dst_path)
                exported_reports[filename] = dst_path

        logger.info(f"Exported {len(exported_reports)} reports to {output_dir}")
        return exported_reports

    def _add_title_slide(self, prs: Presentation, company_name: str, sub_title: str = "Financial Analysis Report"):
        """Add title slide to presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"{company_name}"
        subtitle.text = f"{sub_title}\n{datetime.now().strftime('%B %d, %Y')}"

    def _add_executive_summary_slide(self, prs: Presentation, company_data: Dict[str, Any],
                                   risk_assessment: Dict[str, Any]):
        """Add executive summary slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Executive Summary"

        # Get content placeholder
        content = slide.placeholders[1]

        # Build summary text
        summary_text = self._build_executive_summary_text(company_data, risk_assessment)

        text_frame = content.text_frame
        text_frame.text = summary_text

    def _add_financial_highlights_slide(self, prs: Presentation, company_data: Dict[str, Any]):
        """Add financial highlights slide."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Financial Highlights"

        # Add financial highlights as text boxes
        self._add_financial_highlights_content(slide, company_data)

    def _add_ratio_analysis_slides(self, prs: Presentation, company_data: Dict[str, Any],
                                 charts: Dict[str, str]):
        """Add ratio analysis slides."""
        ratios = company_data.get('analysis', {}).get('ratios', {})

        # Create slides for each ratio category
        for category, ratio_data in ratios.items():
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            title.text = f"{category.title()} Analysis"

            # Add chart if available
            chart_key = f"{category}_chart"
            if chart_key in charts:
                self._add_chart_to_slide(slide, charts[chart_key])

    def _add_trend_analysis_slides(self, prs: Presentation, company_data: Dict[str, Any],
                                 charts: Dict[str, str]):
        """Add trend analysis slides."""
        trends = company_data.get('analysis', {}).get('trends', {})

        # Create slide for trend overview
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Trend Analysis"

        # Add trend chart if available
        if 'trend_analysis' in charts:
            self._add_chart_to_slide(slide, charts['trend_analysis'])

    def _add_risk_assessment_slide(self, prs: Presentation, risk_assessment: Dict[str, Any]):
        """Add risk assessment slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Risk Assessment"

        content = slide.placeholders[1]
        text_frame = content.text_frame

        # Add risk summary
        risk_summary = risk_assessment.get('executive_summary', {})
        risk_text = f"Overall Risk Level: {risk_summary.get('overall_risk_level', 'N/A')}\n"
        risk_text += f"Risk Score: {risk_summary.get('overall_risk_score', 'N/A')}\n"
        risk_text += f"Total Risks: {risk_summary.get('total_risks_identified', 'N/A')}"

        text_frame.text = risk_text

    def _add_valuation_metrics_slide(self, prs: Presentation, company_data: Dict[str, Any]):
        """Add valuation metrics slide."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Valuation Metrics"

        # Add valuation data
        self._add_valuation_content(slide, company_data)

    def _add_investment_recommendations_slide(self, prs: Presentation, company_data: Dict[str, Any],
                                            risk_assessment: Dict[str, Any]):
        """Add investment recommendations slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Investment Recommendations"

        content = slide.placeholders[1]
        text_frame = content.text_frame

        # Generate recommendations
        recommendations = self._generate_investment_recommendations(company_data, risk_assessment)
        text_frame.text = recommendations

    def _add_title_page(self, doc: Document, company_name: str, report_type: str):
        """Add title page to document."""
        # Add title
        title = doc.add_heading(company_name, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add subtitle
        subtitle = doc.add_paragraph(report_type)
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Add date
        date_para = doc.add_paragraph(datetime.now().strftime('%B %d, %Y'))
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_page_break()

    def _add_company_overview_section(self, doc: Document, company_data: Dict[str, Any]):
        """Add company overview section."""
        doc.add_heading('Company Overview', level=1)

        # Add company information
        overview_text = self._build_company_overview_text(company_data)
        doc.add_paragraph(overview_text)

        doc.add_page_break()

    def _add_industry_analysis_section(self, doc: Document, industry_data: Dict[str, Any]):
        """Add industry analysis section."""
        doc.add_heading('Industry Analysis', level=1)

        if industry_data:
            doc.add_paragraph(json.dumps(industry_data, indent=2))
        else:
            doc.add_paragraph("Industry analysis data not available.")

        doc.add_page_break()

    def _add_financial_analysis_section(self, doc: Document, company_data: Dict[str, Any]):
        """Add financial analysis section."""
        doc.add_heading('Financial Analysis', level=1)

        # Add financial highlights
        financial_text = self._build_financial_analysis_text(company_data)
        doc.add_paragraph(financial_text)

        # Add ratios
        doc.add_heading('Key Ratios', level=2)
        ratios = company_data.get('analysis', {}).get('ratios', {})
        for category, ratio_data in ratios.items():
            doc.add_heading(f'{category.title()} Ratios', level=3)
            for key, value in ratio_data.items():
                doc.add_paragraph(f"{key}: {value}")

        doc.add_page_break()

    def _add_competitive_analysis_section(self, doc: Document, competitive_data: Dict[str, Any]):
        """Add competitive analysis section."""
        doc.add_heading('Competitive Analysis', level=1)

        if competitive_data:
            doc.add_paragraph(json.dumps(competitive_data, indent=2))
        else:
            doc.add_paragraph("Competitive analysis data not available.")

        doc.add_page_break()

    def _add_risk_assessment_section(self, doc: Document, risk_data: Dict[str, Any]):
        """Add risk assessment section."""
        doc.add_heading('Risk Assessment', level=1)

        if risk_data:
            doc.add_paragraph(json.dumps(risk_data, indent=2))
        else:
            doc.add_paragraph("Risk assessment data not available.")

        doc.add_page_break()

    def _add_valuation_section(self, doc: Document, company_data: Dict[str, Any],
                             valuation_data: Dict[str, Any]):
        """Add valuation section."""
        doc.add_heading('Valuation Analysis', level=1)

        # Add valuation metrics
        valuation_text = self._build_valuation_text(company_data, valuation_data)
        doc.add_paragraph(valuation_text)

        doc.add_page_break()

    def _add_investment_thesis_section(self, doc: Document, thesis_data: Dict[str, Any]):
        """Add investment thesis section."""
        doc.add_heading('Investment Thesis', level=1)

        if thesis_data:
            doc.add_paragraph(json.dumps(thesis_data, indent=2))
        else:
            doc.add_paragraph("Investment thesis data not available.")

        doc.add_page_break()

    def _add_conclusion_section(self, doc: Document, conclusion_data: Dict[str, Any]):
        """Add conclusion section."""
        doc.add_heading('Conclusion', level=1)

        if conclusion_data:
            doc.add_paragraph(json.dumps(conclusion_data, indent=2))
        else:
            doc.add_paragraph("Conclusion data not available.")

    def _add_key_metrics_slide(self, prs: Presentation, key_metrics: Dict[str, Any]):
        """Add key metrics slide for executive summary."""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Key Metrics"

        # Add key metrics as bullet points
        self._add_key_metrics_content(slide, key_metrics)

    def _add_investment_highlights_slide(self, prs: Presentation, key_metrics: Dict[str, Any]):
        """Add investment highlights slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Investment Highlights"

        content = slide.placeholders[1]
        text_frame = content.text_frame

        highlights_text = self._build_investment_highlights_text(key_metrics)
        text_frame.text = highlights_text

    def _add_risk_summary_slide(self, prs: Presentation, key_metrics: Dict[str, Any]):
        """Add risk summary slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = "Risk Summary"

        content = slide.placeholders[1]
        text_frame = content.text_frame

        risk_text = self._build_risk_summary_text(key_metrics)
        text_frame.text = risk_text

    def _populate_pptx_template(self, template_path: str, data: Dict[str, Any],
                              charts: Dict[str, str]) -> str:
        """Populate PPTX template with data."""
        prs = Presentation(template_path)

        # Update placeholders with data
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    # Replace placeholders with actual data
                    shape.text = self._replace_placeholders(shape.text, data)

        # Add charts to designated areas
        self._insert_charts_to_pptx(prs, charts)

        # Save populated template
        filename = os.path.basename(template_path).replace('.pptx', '_populated.pptx')
        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)

        return output_path

    def _populate_docx_template(self, template_path: str, data: Dict[str, Any],
                              charts: Dict[str, str]) -> str:
        """Populate DOCX template with data."""
        doc = Document(template_path)

        # Replace placeholders in paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text:
                paragraph.text = self._replace_placeholders(paragraph.text, data)

        # Replace placeholders in tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text:
                            paragraph.text = self._replace_placeholders(paragraph.text, data)

        # Save populated template
        filename = os.path.basename(template_path).replace('.docx', '_populated.docx')
        output_path = os.path.join(self.output_dir, filename)
        doc.save(output_path)

        return output_path

    def _build_executive_summary_text(self, company_data: Dict[str, Any],
                                    risk_assessment: Dict[str, Any]) -> str:
        """Build executive summary text."""
        summary = []

        # Financial health
        health = company_data.get('analysis', {}).get('financial_health', {})
        summary.append(f"Financial Health: {health.get('level', 'N/A')} (Score: {health.get('score', 0):.1f})")

        # Key metrics
        ratios = company_data.get('analysis', {}).get('ratios', {})
        if 'profitability' in ratios:
            profitability = ratios['profitability']
            summary.append(f"Profitability: Net margin {profitability.get('net_margin_latest', 0):.1%}")

        # Risk assessment
        risk_summary = risk_assessment.get('executive_summary', {})
        summary.append(f"Risk Level: {risk_summary.get('overall_risk_level', 'N/A')}")

        return "\n".join(summary)

    def _build_company_overview_text(self, company_data: Dict[str, Any]) -> str:
        """Build company overview text."""
        metadata = company_data.get('metadata', {})
        overview = [
            f"Company Name: {metadata.get('company_name', 'N/A')}",
            f"Filing Type: {metadata.get('filing_type', 'N/A')}",
            f"Filing Date: {metadata.get('filing_date', 'N/A')}"
        ]
        return "\n".join(overview)

    def _build_financial_analysis_text(self, company_data: Dict[str, Any]) -> str:
        """Build financial analysis text."""
        analysis = company_data.get('analysis', {})
        trends = analysis.get('trends', {})

        text_parts = []

        for metric, trend_data in trends.items():
            direction = trend_data.get('direction', 'stable')
            change = trend_data.get('percentage_change', 0)
            text_parts.append(f"{metric.title()}: {direction} ({change:+.1f}%)")

        return "\n".join(text_parts)

    def _build_valuation_text(self, company_data: Dict[str, Any],
                            valuation_data: Dict[str, Any]) -> str:
        """Build valuation analysis text."""
        if valuation_data:
            return json.dumps(valuation_data, indent=2)
        else:
            return "Valuation analysis not available."

    def _build_investment_highlights_text(self, key_metrics: Dict[str, Any]) -> str:
        """Build investment highlights text."""
        highlights = []
        for key, value in key_metrics.items():
            highlights.append(f"{key}: {value}")
        return "\n".join(highlights)

    def _build_risk_summary_text(self, key_metrics: Dict[str, Any]) -> str:
        """Build risk summary text."""
        risks = key_metrics.get('risks', {})
        summary = []
        for risk_type, risk_data in risks.items():
            summary.append(f"{risk_type}: {risk_data.get('score', 'N/A')}")
        return "\n".join(summary)

    def _add_financial_highlights_content(self, slide, company_data: Dict[str, Any]):
        """Add financial highlights content to slide."""
        # This would add text boxes with key financial metrics
        # Implementation depends on specific layout requirements
        pass

    def _add_valuation_content(self, slide, company_data: Dict[str, Any]):
        """Add valuation content to slide."""
        # This would add valuation metrics to the slide
        # Implementation depends on specific layout requirements
        pass

    def _add_key_metrics_content(self, slide, key_metrics: Dict[str, Any]):
        """Add key metrics content to slide."""
        # This would add key metrics as bullet points or text boxes
        # Implementation depends on specific layout requirements
        pass

    def _add_chart_to_slide(self, slide, chart_path: str):
        """Add chart image to slide."""
        if os.path.exists(chart_path):
            # Add picture to slide
            left = Inches(1)
            top = Inches(2)
            height = Inches(4.5)
            slide.shapes.add_picture(chart_path, left, top, height=height)

    def _insert_charts_to_pptx(self, prs: Presentation, charts: Dict[str, str]):
        """Insert charts into PPTX template."""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and '{chart:' in shape.text:
                    # Extract chart key from placeholder
                    chart_key = shape.text.split(':')[1].rstrip('}')
                    if chart_key in charts:
                        # Replace text box with chart
                        # This is a simplified implementation
                        pass

    def _replace_placeholders(self, text: str, data: Dict[str, Any]) -> str:
        """Replace placeholders in text with actual data."""
        import re

        # Replace {key} with actual values
        def replace_match(match):
            key = match.group(1)
            return str(data.get(key, match.group(0)))

        return re.sub(r'\{([^}]+)\}', replace_match, text)

    def _generate_investment_recommendations(self, company_data: Dict[str, Any],
                                            risk_assessment: Dict[str, Any]) -> str:
        """Generate investment recommendations."""
        recommendations = []

        # Analyze financial health
        health = company_data.get('analysis', {}).get('financial_health', {})
        risk_level = risk_assessment.get('executive_summary', {}).get('overall_risk_level', 'Medium')

        if health.get('score', 0) >= 80 and risk_level == 'Low':
            recommendations.append("STRONG BUY: Excellent financial health with low risk profile")
        elif health.get('score', 0) >= 60 and risk_level in ['Low', 'Medium']:
            recommendations.append("BUY: Good financial health with manageable risks")
        elif health.get('score', 0) >= 40:
            recommendations.append("HOLD: Fair financial health, monitor risk factors")
        else:
            recommendations.append("SELL: Poor financial health with high risks")

        return "\n".join(recommendations)


def main():
    """Example usage of ReportGenerator."""
    generator = ReportGenerator()

    # Sample data
    sample_company_data = {
        'metadata': {'company_name': 'Sample Company', 'filing_type': '10-K'},
        'analysis': {
            'financial_health': {'score': 75.0, 'level': 'Good'},
            'ratios': {
                'profitability': {'net_margin_latest': 0.12},
                'liquidity': {'current_ratio_latest': 1.5}
            },
            'trends': {
                'revenue': {'direction': 'up', 'percentage_change': 15.5},
                'net_income': {'direction': 'up', 'percentage_change': 20.0}
            }
        }
    }

    sample_risk_assessment = {
        'executive_summary': {
            'overall_risk_level': 'Medium',
            'overall_risk_score': 3.2,
            'total_risks_identified': 8
        }
    }

    sample_charts = {
        'ratio_analysis': 'sample_ratio_chart.png',
        'trend_analysis': 'sample_trend_chart.png'
    }

    # Generate reports
    analysis_report = generator.generate_analysis_report(
        sample_company_data, sample_charts, sample_risk_assessment, "Sample Company"
    )

    investment_thesis = generator.generate_investment_thesis(
        sample_company_data, {}, "Sample Company"
    )

    executive_summary = generator.generate_executive_summary(
        {'revenue': 1000000, 'net_margin': 0.12}, "Sample Company"
    )

    print(f"Generated reports:")
    print(f"- Analysis Report: {analysis_report}")
    print(f"- Investment Thesis: {investment_thesis}")
    print(f"- Executive Summary: {executive_summary}")


if __name__ == "__main__":
    main()